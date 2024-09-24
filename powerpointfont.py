from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

def check_and_update_fonts(pptx_file, target_font):
    presentation = Presentation(pptx_file)
    changes_made = False

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name != target_font:
                            print(f"Found non-matching font: {run.font.name} in slide {presentation.slides.index(slide) + 1}")
                            run.font.name = target_font
                            changes_made = True

    if changes_made:
        presentation.save('updated_presentation.pptx')
        print("Font changes have been applied and saved to 'updated_presentation.pptx'")
    else:
        print("No font changes were necessary")

# Usage
pptx_file = 'path/to/your/presentation.pptx'
target_font = 'Arial'  # Specify your desired font here
check_and_update_fonts(pptx_file, target_font)