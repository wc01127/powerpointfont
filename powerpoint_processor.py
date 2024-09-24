from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def process_powerpoint(input_file, output_file, target_font):
    presentation = Presentation(input_file)
    total_text_runs = 0
    changed_text_runs = 0
    slides_changed = set()

    print(f"Target font: {target_font}")

    for slide_number, slide in enumerate(presentation.slides, start=1):
        print(f"\nProcessing Slide {slide_number}")
        slide_changed = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        total_text_runs += 1
                        print(f"  Text: '{run.text}', Current font: '{run.font.name}'")
                        
                        # Check if the font is not set or different from the target font
                        if not run.font.name or run.font.name.lower() != target_font.lower():
                            print(f"    Changing font to {target_font}")
                            run.font.name = target_font
                            changed_text_runs += 1
                            slide_changed = True
                        else:
                            print(f"    Font already matches {target_font}")
        
        if slide_changed:
            slides_changed.add(slide_number)
        print(f"Slide {slide_number} changed: {slide_changed}")

    presentation.save(output_file)
    
    percent_changed = (changed_text_runs / total_text_runs * 100) if total_text_runs > 0 else 0
    print(f"\nTotal text runs: {total_text_runs}")
    print(f"Changed text runs: {changed_text_runs}")
    print(f"Percent changed: {percent_changed:.2f}%")
    print(f"Slides changed: {slides_changed}")
    
    return percent_changed, slides_changed