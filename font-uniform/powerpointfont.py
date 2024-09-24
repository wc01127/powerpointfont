import os
import sys
from pptx import Presentation
from collections import Counter


def analyze_fonts(input_pptx):
    prs = Presentation(input_pptx)
    print(f"Analyzing fonts in '{os.path.basename(input_pptx)}':")

    for slide_index, slide in enumerate(prs.slides, 1):
        fonts = []
        print(f"\nSlide {slide_index}:")

        # Combine placeholders and shapes into a single list
        all_shapes = list(slide.placeholders) + \
            [shape for shape in slide.shapes if shape not in slide.placeholders]

        for shape_index, shape in enumerate(all_shapes, 1):
            shape_fonts = []
            shape_text = ""

            if shape.has_text_frame:
                shape_text = shape.text
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            shape_fonts.append(run.font.name)
                        elif run.font.size:
                            shape_fonts.append(
                                f"Unknown font (size: {run.font.size})")
                        else:
                            shape_fonts.append("Default font")

            print(f"  Shape {shape_index} ({shape.name}):")
            print(f"    Text: {shape_text}")
            if shape_fonts:
                fonts.extend(shape_fonts)
                print(f"    Fonts: {', '.join(set(shape_fonts))}")
            else:
                print("    No fonts detected")

            # Additional shape information
            print(f"    Shape type: {type(shape)}")
            if hasattr(shape, 'shape_type'):
                print(f"    Shape type: {shape.shape_type}")
            if hasattr(shape, 'placeholder_format'):
                print(f"    Placeholder type: {shape.placeholder_format.type}")

        font_count = Counter(fonts)
        if font_count:
            print("  Summary:")
            for font, count in font_count.items():
                print(
                    f"    - {font}: {count} occurrence{'s' if count > 1 else ''}")
        else:
            print("  No fonts detected on this slide.")

        print(f"  Total shapes on slide: {len(all_shapes)}")


def set_uniform_font(input_pptx, output_pptx, target_font_name):
    # Load the presentation
    prs = Presentation(input_pptx)
    changes_made = False

    # Iterate through slides
    for slide in prs.slides:
        # Iterate through shapes
        for shape in slide.shapes:
            # Handle shapes with text
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name != target_font_name:
                            run.font.name = target_font_name
                            changes_made = True
            # Handle tables
            elif shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.name != target_font_name:
                                    run.font.name = target_font_name
                                    changes_made = True
            # Handle charts
            elif shape.has_chart:
                chart = shape.chart
                # Update chart title
                if chart.has_title:
                    for paragraph in chart.chart_title.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.name != target_font_name:
                                run.font.name = target_font_name
                                changes_made = True
                # Update axis titles and labels
                for axis in (chart.category_axis, chart.value_axis):
                    if axis.has_title:
                        for paragraph in axis.axis_title.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.name != target_font_name:
                                    run.font.name = target_font_name
                                    changes_made = True
                    # Update axis tick labels (if applicable)
                    for tick_label in axis.tick_labels:
                        if tick_label.font.name != target_font_name:
                            tick_label.font.name = target_font_name
                            changes_made = True
                # Update data labels
                for series in chart.series:
                    if series.has_data_labels:
                        for data_label in series.data_labels:
                            if data_label.font.name != target_font_name:
                                data_label.font.name = target_font_name
                                changes_made = True

    # Save the updated presentation only if changes were made
    if changes_made:
        prs.save(output_pptx)
        print(
            f"Fonts updated to '{target_font_name}'. Saved as '{output_pptx}'.")
    else:
        print(
            f"No font changes were necessary. All fonts are already '{target_font_name}'.")

    return changes_made


# New function to list and choose test slides
def choose_test_slide():
    test_slides_folder = "test_slides"
    slides = [f for f in os.listdir(test_slides_folder) if f.endswith(
        ".pptx") and not f.endswith("_updated.pptx")]

    if not slides:
        print("No test slides found in the 'test_slides' folder.")
        return None

    print("Available test slides:")
    for i, slide in enumerate(slides, 1):
        print(f"{i}. {slide}")

    while True:
        try:
            choice = int(
                input("Enter the number of the slide you want to use: "))
            if 1 <= choice <= len(slides):
                return os.path.join(test_slides_folder, slides[choice - 1])
            else:
                print("Invalid choice. Please try again.")
        except ValueError:
            print("Please enter a valid number.")


# Move the main execution code into a separate function
def main():
    input_pptx = choose_test_slide()
    if input_pptx:
        analyze_fonts(input_pptx)

        proceed = input(
            "\nDo you want to proceed with font uniformization? (y/n): ").lower()
        if proceed == 'y':
            base_name = os.path.splitext(os.path.basename(input_pptx))[0]
            output_pptx = os.path.join(
                "test_slides", f"{base_name}_updated.pptx")
            target_font_name = input("Enter the desired font name: ")
            set_uniform_font(input_pptx, output_pptx, target_font_name)
        else:
            print("Font uniformization cancelled.")
    else:
        print("No input file selected. Exiting.")


# Use this idiom to make the script both importable and executable
if __name__ == "__main__":
    if len(sys.argv) != 4:
        print(
            "Usage: python3 powerpointfont.py <input_pptx> <output_pptx> <target_font_name>")
        sys.exit(1)

    input_pptx = sys.argv[1]
    output_pptx = sys.argv[2]
    target_font_name = sys.argv[3]

    analyze_fonts(input_pptx)
    set_uniform_font(input_pptx, output_pptx, target_font_name)
